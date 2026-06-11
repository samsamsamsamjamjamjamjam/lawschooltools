#!/usr/bin/env python3
"""
Law School Tools — PPT generation backend.

Opens template1.pptx, fills slide placeholders with the user's notes,
returns a finished .pptx file. No AI — pure rule-based polish.
"""
import copy
import io
import os
import re
from pathlib import Path

from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Pt

HERE = Path(__file__).parent
# Template 3 — cleanest/simplest (Title Slide, Title and Content, Section Header, etc.)
TEMPLATE_PATH = HERE / "template3.pptx"

app = FastAPI()
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


class GenerateRequest(BaseModel):
    notes: str
    title: str = "Study Deck"
    presenter: str = ""
    polish: bool = True
    length: str = "standard"  # "short", "standard", "detailed"


# -----------------------------
# Rule-based polish (no AI)
# -----------------------------
FILLER_WORDS = {
    "basically", "literally", "actually", "honestly", "really", "very",
    "just", "stuff", "things", "kind of", "sort of", "like,",
}


def polish_text(text: str) -> str:
    """Cleanup notes: trim filler, capitalize, ensure terminal punctuation."""
    if not text:
        return text
    t = text.strip()
    # Strip common filler words (case-insensitive, word-boundary)
    for filler in FILLER_WORDS:
        pattern = r"\b" + re.escape(filler) + r"\b\s*"
        t = re.sub(pattern, "", t, flags=re.IGNORECASE)
    # Collapse whitespace
    t = re.sub(r"\s+", " ", t).strip()
    # Capitalize first letter of each sentence
    sentences = re.split(r"(?<=[.!?])\s+", t)
    sentences = [s[0].upper() + s[1:] if s and s[0].isalpha() else s for s in sentences]
    t = " ".join(sentences).strip()
    # Ensure terminal punctuation
    if t and t[-1] not in ".!?":
        t += "."
    return t


# -----------------------------
# Note extraction
# -----------------------------
def extract_concepts(notes: str):
    """
    Pull (concept, definition) pairs from notes.
    Looks for patterns like "X is Y", "X means Y", "X requires Y", "X: Y".
    Falls back to chunked sentences if no patterns match.
    """
    sentences = re.split(r"(?<=[.!?])\s+", notes.replace("\n", " ").strip())
    sentences = [s.strip() for s in sentences if len(s.strip()) > 12]

    patterns = [
        # "The X doctrine/rule/test/standard ... Y" — match first so it wins over generic "X is Y"
        r"^(?:The\s+)?(.+?\s+(?:doctrine|rule|test|standard))\s+(?:is|requires|states(?:\s+that)?|means)\s+(.+)$",
        r"^(.+?)\s+(?:is|are|means|refers to)\s+(.+)$",
        r"^(.+?)\s+requires?\s+(.+)$",
        r"^(.+?):\s+(.+)$",
    ]

    facts = []
    for s in sentences:
        for pat in patterns:
            m = re.match(pat, s, flags=re.IGNORECASE)
            if m:
                concept = m.group(1).strip().lstrip("Tt").lstrip("he ").strip()
                # better: strip leading "The "
                concept = re.sub(r"^the\s+", "", m.group(1).strip(), flags=re.IGNORECASE)
                definition = m.group(2).strip().rstrip(".!?")
                if 2 < len(concept) < 80 and 5 < len(definition) < 300:
                    facts.append({"concept": concept, "definition": definition, "sentence": s})
                    break

    if not facts:
        # Fallback: chunk sentences into pseudo-facts
        for i, s in enumerate(sentences[:15]):
            facts.append({"concept": f"Key Point {i+1}", "definition": s, "sentence": s})

    return facts


# -----------------------------
# Template filling
# -----------------------------
def set_text_preserving_format(text_frame, new_text: str):
    """
    Replace text in a text_frame while keeping the first run's formatting.
    pptx text frames have paragraphs -> runs. Wiping and rewriting loses
    fonts, sizes, colors. So we keep the first run and replace its text.
    """
    if not text_frame.paragraphs:
        text_frame.text = new_text
        return

    # Keep first paragraph + first run; clear the rest
    first_para = text_frame.paragraphs[0]
    if not first_para.runs:
        first_para.text = new_text
        return

    first_run = first_para.runs[0]
    first_run.text = new_text

    # Remove additional runs in first paragraph
    for run in first_para.runs[1:]:
        run._r.getparent().remove(run._r)

    # Remove additional paragraphs
    for para in text_frame.paragraphs[1:]:
        para._p.getparent().remove(para._p)


def find_placeholder_by_text(slide, search_text: str):
    """Find a shape whose text contains search_text (case-insensitive)."""
    search_lower = search_text.lower()
    for shape in slide.shapes:
        if shape.has_text_frame and search_lower in shape.text_frame.text.lower():
            return shape
    return None


def get_title_placeholder(slide):
    """Find the title placeholder (largest text shape near top, or PP_PLACEHOLDER.TITLE)."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    for shape in slide.placeholders:
        try:
            if shape.placeholder_format.idx == 0:  # Title is always idx 0
                return shape
        except Exception:
            continue
    # Fallback: largest text shape near top
    candidates = [s for s in slide.shapes if s.has_text_frame and s.top is not None and s.top < 2000000]
    if candidates:
        return max(candidates, key=lambda s: (s.width or 0) * (s.height or 0))
    return None


def get_body_placeholder(slide):
    """Find the main body/content placeholder."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    body_idx = None
    body_shape = None
    largest_area = 0
    for shape in slide.placeholders:
        try:
            idx = shape.placeholder_format.idx
            if idx == 0:
                continue  # skip title
            if shape.has_text_frame:
                area = (shape.width or 0) * (shape.height or 0)
                if area > largest_area:
                    largest_area = area
                    body_shape = shape
        except Exception:
            continue
    return body_shape


def _drop_all_slides(prs):
    """Properly remove every slide from a Presentation — both the XML refs
    in the slide list AND the underlying slide parts/relationships,
    so the saved .pptx isn't bloated with orphan slides."""
    sldIdLst = prs.slides._sldIdLst
    # Collect rIds before mutation
    rels_to_drop = []
    for sldId in list(sldIdLst):
        rels_to_drop.append(sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'))
        sldIdLst.remove(sldId)
    part = prs.part
    for rId in rels_to_drop:
        if rId and rId in part.rels:
            try:
                part.drop_rel(rId)
            except Exception:
                pass


def add_slide_from_layout(prs, layout_name: str):
    """Add a new slide using a named layout. Returns the slide."""
    for layout in prs.slide_layouts:
        if layout.name == layout_name:
            return prs.slides.add_slide(layout)
    # Fallback: first content layout
    return prs.slides.add_slide(prs.slide_layouts[3])  # "Title and Content"


# -----------------------------
# Main generator
# -----------------------------
def generate_deck(notes: str, title: str, presenter: str, polish: bool, length: str) -> bytes:
    """Build a filled .pptx and return its bytes."""
    if not TEMPLATE_PATH.exists():
        raise HTTPException(500, f"Template not found at {TEMPLATE_PATH}")

    facts = extract_concepts(notes)
    if not facts:
        raise HTTPException(400, "Could not extract any concepts from the notes. Try adding sentences like 'X is Y' or 'X means Y'.")

    # Decide deck length
    length_map = {"short": 5, "standard": 9, "detailed": 14}
    target_content_slides = min(length_map.get(length, 9), len(facts))

    # Apply polish
    if polish:
        for f in facts:
            f["concept"] = polish_text(f["concept"]).rstrip(".")
            f["definition"] = polish_text(f["definition"])
            title = polish_text(title).rstrip(".")

    # Open template fresh
    prs = Presentation(str(TEMPLATE_PATH))

    # Strategy: properly drop all existing example slides (remove rels + parts),
    # then add fresh slides from the layouts we want.
    _drop_all_slides(prs)

    # --- 1. Title slide ---
    title_slide = add_slide_from_layout(prs, "Title Slide")
    title_ph = get_title_placeholder(title_slide)
    if title_ph and title_ph.has_text_frame:
        set_text_preserving_format(title_ph.text_frame, title)
    # presenter line
    for shape in title_slide.placeholders:
        if shape.has_text_frame and shape != title_ph:
            txt = shape.text_frame.text.lower()
            if "presented by" in txt or "insert name" in txt:
                set_text_preserving_format(
                    shape.text_frame,
                    f"Presented by {presenter}" if presenter else "Law School Tools",
                )
                break

    # --- 2. Section Header (Introduction / Overview) ---
    section = add_slide_from_layout(prs, "Section Header")
    section_title = get_title_placeholder(section)
    if section_title and section_title.has_text_frame:
        set_text_preserving_format(section_title.text_frame, "Overview")
    # subtitle/body if present
    section_body = get_body_placeholder(section)
    if section_body and section_body.has_text_frame:
        overview = f"This presentation covers {len(facts[:target_content_slides])} key concepts."
        set_text_preserving_format(section_body.text_frame, overview)

    # --- 3..N. Content slides — one concept each ---
    for fact in facts[:target_content_slides]:
        slide = add_slide_from_layout(prs, "Title and Content")
        t_ph = get_title_placeholder(slide)
        if t_ph and t_ph.has_text_frame:
            set_text_preserving_format(t_ph.text_frame, fact["concept"])
        b_ph = get_body_placeholder(slide)
        if b_ph and b_ph.has_text_frame:
            set_text_preserving_format(b_ph.text_frame, fact["definition"])

    # --- Final. Conclusion slide (template3 has no "Conclusion" layout, so reuse Title and Content) ---
    conclusion = add_slide_from_layout(prs, "Title and Content")
    c_title = get_title_placeholder(conclusion)
    if c_title and c_title.has_text_frame:
        set_text_preserving_format(c_title.text_frame, "Conclusion")
    c_body = get_body_placeholder(conclusion)
    if c_body and c_body.has_text_frame:
        # Build a short recap
        recap_items = [polish_text(f["concept"]).rstrip(".") for f in facts[:target_content_slides]]
        recap = "Key takeaways: " + "; ".join(recap_items) + "."
        set_text_preserving_format(c_body.text_frame, recap)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# -----------------------------
# Routes
# -----------------------------
@app.get("/api/health")
def health():
    return {"ok": True, "template_present": TEMPLATE_PATH.exists()}


@app.post("/api/generate-ppt")
def generate_ppt(req: GenerateRequest):
    if not req.notes.strip():
        raise HTTPException(400, "Notes are required.")
    data = generate_deck(
        notes=req.notes,
        title=req.title or "Study Deck",
        presenter=req.presenter or "",
        polish=req.polish,
        length=req.length or "standard",
    )
    safe_name = re.sub(r"[^a-zA-Z0-9]+", "-", req.title or "study-deck").strip("-").lower() or "study-deck"
    return StreamingResponse(
        io.BytesIO(data),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{safe_name}.pptx"'},
    )


if __name__ == "__main__":
    import os
    import uvicorn
    port = int(os.environ.get("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
